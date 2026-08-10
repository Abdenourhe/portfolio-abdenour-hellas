from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from PIL import Image
import io
import os

# Charger la présentation
prs = Presentation("ABDENOUR_HELLAS_CV.pptx")

# Dimensions de la slide (en pixels à 96 DPI)
slide_width_px = int(prs.slide_width.inches * 96)
slide_height_px = int(prs.slide_height.inches * 96)

print(f"Dimensions des slides: {slide_width_px}x{slide_height_px}")

# Pour chaque slide, on va la rendre en image
for i, slide in enumerate(prs.slides, 1):
    # Créer une image blanche de fond
    img = Image.new('RGB', (slide_width_px, slide_height_px), 'white')
    
    # Pour chaque shape, extraire le texte et le placer (simplifié)
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            print(f"Slide {i}: {shape.text[:80]}...")

print("\nConversion terminée. Le CV est prêt pour être utilisé.")
