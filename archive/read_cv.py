from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation("ABDENOUR_HELLAS_CV.pptx")

print("=== CONTENU DU CV ===\n")

for i, slide in enumerate(prs.slides, 1):
    print(f"\n--- SLIDE {i} ---")
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            print(shape.text.strip())
            print("-" * 40)
