from pptx import Presentation
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
from reportlab.lib import colors
from reportlab.lib.enums import TA_LEFT, TA_CENTER
import os

# Lire le PPT
prs = Presentation("ABDENOUR_HELLAS_CV.pptx")

# Créer le PDF
pdf_path = "ABDENOUR_HELLAS_CV.pdf"
doc = SimpleDocTemplate(pdf_path, pagesize=letter,
                        rightMargin=0.6*inch, leftMargin=0.6*inch,
                        topMargin=0.6*inch, bottomMargin=0.6*inch)

styles = getSampleStyleSheet()

# Styles personnalisés
title_style = ParagraphStyle(
    'CustomTitle',
    parent=styles['Heading1'],
    fontSize=20,
    textColor=colors.HexColor('#1a365d'),
    spaceAfter=6,
    alignment=TA_CENTER,
    fontName='Helvetica-Bold'
)

subtitle_style = ParagraphStyle(
    'CustomSubtitle',
    parent=styles['Normal'],
    fontSize=11,
    textColor=colors.HexColor('#4a5568'),
    alignment=TA_CENTER,
    spaceAfter=12,
    fontName='Helvetica'
)

heading_style = ParagraphStyle(
    'CustomHeading',
    parent=styles['Heading2'],
    fontSize=13,
    textColor=colors.HexColor('#1a365d'),
    spaceAfter=6,
    spaceBefore=12,
    fontName='Helvetica-Bold',
    borderWidth=0,
    borderColor=colors.HexColor('#1a365d'),
    borderPadding=4,
    leftIndent=0,
    backColor=colors.HexColor('#e2e8f0')
)

normal_style = ParagraphStyle(
    'CustomNormal',
    parent=styles['Normal'],
    fontSize=10,
    leading=14,
    spaceAfter=4,
    fontName='Helvetica'
)

# Extraire le texte par slide
def extract_text_from_slide(slide):
    texts = []
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            texts.append(shape.text.strip())
    return texts

story = []

# === SLIDE 1 - PROFIL ET EXPÉRIENCE ===
texts1 = extract_text_from_slide(prs.slides[0])

# En-tête
story.append(Paragraph("ABDENOUR HELLAS", title_style))
story.append(Paragraph("SPÉCIALISTE EN GÉNIE ÉLECTRIQUE & DÉVELOPPEMENT WEB", subtitle_style))
story.append(Paragraph("Nouveau-Brunswick, Canada | abdenour.hellas@uqat.ca | +1 418 350 5686 | https://abdenour-hellas.online", subtitle_style))
story.append(Spacer(1, 12))

# Profil
story.append(Paragraph("PROFIL", heading_style))
story.append(Paragraph("Spécialiste en génie électrique avec double compétence en développement web, orienté vers la conception de systèmes électriques collaboratifs et l'automatisation industrielle. 5+ ans d'expérience internationale (Algérie, Canada) en maintenance, supervision et conception. Passionné par l'innovation énergétique et les solutions numériques. Recherche des opportunités à l'international.", normal_style))
story.append(Spacer(1, 6))

# Expérience
story.append(Paragraph("EXPÉRIENCE PROFESSIONNELLE", heading_style))

exp_data = [
    [Paragraph("<b>Technicien en Génie Électrique</b> — ENTP, In-Amenas, Algérie", normal_style), Paragraph("déc 2025 - mars 2026", normal_style)],
    [Paragraph("• Conception et maintenance du système collaboratif électrique TCOST (normes, standards, options de configuration)<br/>• Développement d'une bibliothèque numérique complète : produits, prix, spécifications techniques, codes applicables<br/>• Intégration de solutions innovantes : systèmes intelligents, efficacité énergétique, automatisations", normal_style), ""],
    [Paragraph("<b>Superviseur Département Électroménagers</b> — RONA, Rouyn-Noranda, Canada", normal_style), Paragraph("juin 2023 - nov 2025", normal_style)],
    [Paragraph("• Supervision des opérations du département électroménagers<br/>• Gestion d'équipe et optimisation des processus", normal_style), ""],
    [Paragraph("<b>Électricien Industriel</b> — Maisons Laprise, Québec, Canada", normal_style), Paragraph("nov 2018 - août 2019", normal_style)],
    [Paragraph("• Montage, réparation et entretien des circuits électriques, tableaux de distribution et dispositifs de commande<br/>• Diagnostic et résolution des dysfonctionnements (surchauffe, courts-circuits, défauts d'isolation)<br/>• Rédaction des rapports d'intervention et mise à jour des schémas électriques", normal_style), ""],
    [Paragraph("<b>Stages Techniques</b> — GTP / SNP / Complexe Industriel TDA, Algérie", normal_style), Paragraph("jan 2017 - déc 2019", normal_style)],
    [Paragraph("• Maintenance préventive/corrective des équipements électriques et industriels<br/>• Inspection qualité, supervision travaux, gestion des stocks de pièces de rechange", normal_style), ""],
    [Paragraph("<b>Professeur d'enseignement secondaire</b> — Lycée Sheikh Amoud, In-amenas, Algérie", normal_style), Paragraph("fév 2022 - août 2022", normal_style)],
    [Paragraph("• Évaluation des progrès des élèves, retour constructif pour accompagner la progression", normal_style), ""],
]

exp_table = Table(exp_data, colWidths=[doc.width*0.75, doc.width*0.25])
exp_table.setStyle(TableStyle([
    ('VALIGN', (0,0), (-1,-1), 'TOP'),
    ('LEFTPADDING', (0,0), (-1,-1), 4),
    ('RIGHTPADDING', (0,0), (-1,-1), 4),
    ('BOTTOMPADDING', (0,0), (-1,-1), 6),
]))
story.append(exp_table)

# === SLIDE 2 - ÉDUCATION, COMPÉTENCES, PROJETS ===
texts2 = extract_text_from_slide(prs.slides[1])

story.append(Paragraph("ÉDUCATION", heading_style))
edu_data = [
    [Paragraph("<b>M.Eng. Génie Électrique</b> — UQAT, Rouyn-Noranda, Canada", normal_style), Paragraph("juin 2025", normal_style)],
    [Paragraph("<b>Master Électromécanique</b> — Université de Batna2, Algérie", normal_style), Paragraph("juin 2020", normal_style)],
    [Paragraph("<b>Licence Électromécanique</b> — Université de Batna2, Algérie", normal_style), Paragraph("juin 2018", normal_style)],
    [Paragraph("<b>Technicien Supérieur Informatique-Réseaux</b> — CFPA, Algérie", normal_style), Paragraph("sept 2022", normal_style)],
    [Paragraph("<b>Baccalauréat en Science Expérimentale</b> — Lycée Kser Belezma, Batna, Algérie", normal_style), Paragraph("juin 2015", normal_style)],
]
edu_table = Table(edu_data, colWidths=[doc.width*0.75, doc.width*0.25])
edu_table.setStyle(TableStyle([
    ('VALIGN', (0,0), (-1,-1), 'TOP'),
    ('LEFTPADDING', (0,0), (-1,-1), 4),
    ('BOTTOMPADDING', (0,0), (-1,-1), 4),
]))
story.append(edu_table)

story.append(Paragraph("CERTIFICATS", heading_style))
story.append(Paragraph("• ASP Construction (Santé & Sécurité) — Centre Polymétier, Rouyn-Noranda, Canada (juill 2024 - sept 2024)", normal_style))
story.append(Paragraph("• <b>Admission en cours à l'Ordre des ingénieurs du Québec</b> — Processus de Qualification Professionnelle Initié", normal_style))

story.append(Paragraph("COMPÉTENCES", heading_style))
skills_data = [
    [Paragraph("<b>Électrique</b>", normal_style), Paragraph("Génie Électrique, Électromécanique, Maintenance Industrielle, Instrumentation & Systèmes Embarqués, Automatisation, Câblage Industriel", normal_style)],
    [Paragraph("<b>Web</b>", normal_style), Paragraph("Next.js, React, TypeScript, Prisma, PostgreSQL, Tailwind CSS, Node.js, GitHub", normal_style)],
    [Paragraph("<b>Logiciel</b>", normal_style), Paragraph("Proteus 8, SolidWorks, MATLAB, AutoCAD, Excel avancé", normal_style)],
    [Paragraph("<b>Normes</b>", normal_style), Paragraph("Code québécois de l'électricité (CQE), Canadian Electrical Code (CEC), NFPA 70", normal_style)],
    [Paragraph("<b>Soft</b>", normal_style), Paragraph("Autonomie, Leadership, Rigueur, Travail d'équipe, Gestion de projet", normal_style)],
]
skills_table = Table(skills_data, colWidths=[doc.width*0.18, doc.width*0.82])
skills_table.setStyle(TableStyle([
    ('VALIGN', (0,0), (-1,-1), 'TOP'),
    ('LEFTPADDING', (0,0), (-1,-1), 4),
    ('BOTTOMPADDING', (0,0), (-1,-1), 4),
    ('BACKGROUND', (0,0), (0,-1), colors.HexColor('#e2e8f0')),
]))
story.append(skills_table)

story.append(Paragraph("PROJETS", heading_style))
story.append(Paragraph("• <b>TAHFIDZ</b> — Plateforme SaaS de gestion d'écoles coraniques (Next.js, PostgreSQL, Prisma) — https://tahfidz-two.vercel.app", normal_style))
story.append(Paragraph("• <b>CCI Montmagny</b> — Site web communautaire pour le Centre Culturel Islamique de Montmagny — https://mosquee-cci-montmagny.vercel.app", normal_style))
story.append(Paragraph("• <b>Radar Avancé</b> — Projet de fin d'études UQAT : Système intelligent de surveillance minère", normal_style))
story.append(Paragraph("• <b>Portfolio</b> — https://abdenour-hellas.online", normal_style))

story.append(Paragraph("LANGUES", heading_style))
story.append(Paragraph("Arabe — Natif | Français — Courant | Anglais — Professionnel", normal_style))

story.append(Paragraph("RÉFÉRENCES", heading_style))
story.append(Paragraph("LinkedIn : https://www.linkedin.com/in/abdenour-hellas/", normal_style))

# Construire le PDF
doc.build(story)
print(f"PDF créé avec succès : {os.path.abspath(pdf_path)} ({os.path.getsize(pdf_path)} octets)")
